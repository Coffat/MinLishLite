import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
# Set to 16:9 widescreen
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(8, 51, 112) # #083370
LIGHT_BLUE = RGBColor(112, 182, 229) # #70B6E5
PINK = RGBColor(243, 217, 230) # #F3D9E6
WHITE = RGBColor(255, 255, 255)
BG_COLOR = RGBColor(248, 249, 250) # #F8F9FA
GREY = RGBColor(51, 65, 85)

# Create a blank slide layout
blank_slide_layout = prs.slide_layouts[6]

def set_shape_color(shape, fill_color, line_color=None, line_width=None):
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    if line_width:
        shape.line.width = Pt(line_width)

def add_title(slide, text, left, top, width, height, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = align
    
    # Add a decorative underline
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Pt(55), width/3, Pt(3)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = DARK_BLUE
    slide.shapes[-1].line.fill.background()

def add_bullet_points(slide, points, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, point in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.name = 'Arial'
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_BLUE
        p.level = 0
        p.space_after = Pt(14)

def add_corner_blob(slide, is_top_left=True, color=LIGHT_BLUE):
    if is_top_left:
        blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    else:
        blob = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.333), Inches(3.5), Inches(6), Inches(6))
    set_shape_color(blob, color)

def add_phone_mockup(slide, left, top, title="SCREENSHOT"):
    # Phone border
    phone = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.2), Inches(5.8))
    set_shape_color(phone, WHITE, DARK_BLUE, 4)
    
    # Phone screen (Image Placeholder)
    screen = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.15), top + Inches(0.15), Inches(2.9), Inches(5.5))
    set_shape_color(screen, GREY)
    
    # Text in screen
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(2.5), Inches(2.9), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = f"{title}\nINSERT HERE"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# --- Slide 1: Study / Flashcards ---
slide1 = prs.slides.add_slide(blank_slide_layout)
slide1.background.fill.solid()
slide1.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide1, False, LIGHT_BLUE)
add_title(slide1, "Study / Flashcards", Inches(1), Inches(1), Inches(6), Inches(1))
add_bullet_points(slide1, [
    "Spaced Repetition System (SRS): Custom SrsCalculator schedules reviews based on Again, Hard, Good, and Easy ratings.",
    "Study Modes: Filter by Specific Deck (DeckDue) or Review Today (DueToday) for all overdue words.",
    "Real-time Progress: Live tracking of study fractions and remaining cards within a session."
], Inches(1), Inches(2.5), Inches(6.5), Inches(4))
add_phone_mockup(slide1, Inches(8.5), Inches(1), "STUDY")

# --- Slide 2: Progress & Analytics ---
slide2 = prs.slides.add_slide(blank_slide_layout)
slide2.background.fill.solid()
slide2.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide2, False, LIGHT_BLUE)
add_phone_mockup(slide2, Inches(1.5), Inches(1), "PROGRESS")
add_title(slide2, "Progress & Analytics", Inches(5.5), Inches(1), Inches(6.5), Inches(1))
add_bullet_points(slide2, [
    "Continuous Tracking: Monitors current streak, longest streak, and total reviews mapped from ReviewHistoryEntity.",
    "Retention Insights: Automatically computes retention rate and accuracy metrics.",
    "State Management: Driven by Coroutine Flows delivering real-time ProgressAnalytics updates to the UI."
], Inches(5.5), Inches(2.5), Inches(6.5), Inches(4))

# --- Slide 3: Settings & Profile ---
slide3 = prs.slides.add_slide(blank_slide_layout)
slide3.background.fill.solid()
slide3.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide3, False, LIGHT_BLUE)
add_title(slide3, "Settings & Profile", Inches(1), Inches(1), Inches(6.5), Inches(1))
add_bullet_points(slide3, [
    "Study Goals: Personalize daily target for new words (default 10).",
    "Notifications: Configure customizable daily review reminders (default 09:00 AM).",
    "Profile Management: Edit profile attributes (Name, Goal, Level) synced deeply with UserEntity."
], Inches(1), Inches(2.5), Inches(6.5), Inches(4))
add_phone_mockup(slide3, Inches(8.5), Inches(1), "SETTINGS")

# --- Slide 4: Core Data Models ---
slide4 = prs.slides.add_slide(blank_slide_layout)
slide4.background.fill.solid()
slide4.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide4, True, LIGHT_BLUE)
slide4.shapes.add_picture('datacore.png', Inches(1.0), Inches(1.5), width=Inches(4.5))
add_title(slide4, "Core Data Models", Inches(6.0), Inches(1), Inches(6.5), Inches(1))
add_bullet_points(slide4, [
    "DeckEntity: Represents a vocabulary deck holding multiple flashcards.",
    "WordEntity: Represents a flashcard mapping word data to critical SRS attributes (easeFactor, nextReviewAt, reviewCount, correctCount).",
    "ReviewHistoryEntity: Logs every review session outcome to fuel analytics and progress reporting.",
    "UserEntity: Stores the user's base profile and study goals locally."
], Inches(6.0), Inches(2.5), Inches(6.5), Inches(4))

# --- Slide 5: Navigation Flow ---
slide5 = prs.slides.add_slide(blank_slide_layout)
slide5.background.fill.solid()
slide5.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide5, False, LIGHT_BLUE)
add_title(slide5, "Navigation Flow", Inches(1), Inches(1), Inches(6.5), Inches(1))
add_bullet_points(slide5, [
    "Splash Screen → Home Screen",
    "From the Home Screen navigates to:",
    "  • Deck List",
    "  • Deck Detail",
    "  • Study Mode / Flashcard Review",
    "  • Progress & Analytics",
    "  • Settings & Profile Management"
], Inches(1), Inches(2.5), Inches(6.5), Inches(4))

slide5.shapes.add_picture('navflow.png', Inches(8.0), Inches(1.5), width=Inches(4.5))

# --- Slide 6: Algorithm & Logic ---
slide6 = prs.slides.add_slide(blank_slide_layout)
slide6.background.fill.solid()
slide6.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide6, True, LIGHT_BLUE)
slide6.shapes.add_picture('algorithm.png', Inches(1.0), Inches(1.5), width=Inches(4.5))
add_title(slide6, "Algorithm & Logic", Inches(6.0), Inches(1), Inches(6.5), Inches(1))
add_bullet_points(slide6, [
    "Spaced Repetition System (SRS): SrsCalculator computes the next review time (nextReviewAt) and modifies difficulty (easeFactor) relying on user self-evaluation (Again, Hard, Good, Easy).",
    "Progress Calculation: ProgressCalculator derives learning streaks, tracking accuracy and review retention metrics extracted from ReviewHistoryEntity entries."
], Inches(6.0), Inches(2.5), Inches(6.5), Inches(4))

# --- Slide 7: Current Strengths ---
slide7 = prs.slides.add_slide(blank_slide_layout)
slide7.background.fill.solid()
slide7.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide7, False, LIGHT_BLUE)
add_title(slide7, "Current Strengths", Inches(1), Inches(1), Inches(11.3), Inches(1))
add_bullet_points(slide7, [
    "Modern UI implementation leveraging Jetpack Compose and Material 3.",
    "Robust local-first architecture securely powered by Room Database.",
    "Highly effective custom Spaced Repetition (SRS) algorithmic integration.",
    "Comprehensive analytics and progress tracking systems natively implemented.",
    "Clean MVVM architecture efficiently decoupling UI elements from business logic."
], Inches(1), Inches(2.5), Inches(11.3), Inches(4))

# --- Slide 8: Current Limitations ---
slide8 = prs.slides.add_slide(blank_slide_layout)
slide8.background.fill.solid()
slide8.background.fill.fore_color.rgb = BG_COLOR
add_corner_blob(slide8, True, PINK)
add_title(slide8, "Current Limitations", Inches(1), Inches(1), Inches(11.3), Inches(1))
add_bullet_points(slide8, [
    "Relies on manual Dependency Injection (AppContainer) rather than standard DI frameworks like Hilt or Koin.",
    "Limited cloud data synchronization capabilities (Data currently acts primarily as offline storage).",
    "Heavy dependency on public and free translation/dictionary APIs which are susceptible to rate limits."
], Inches(1), Inches(2.5), Inches(11.3), Inches(4))

prs.save('MinLishLite_Presentation.pptx')
print("Successfully created MinLishLite_Presentation.pptx")
